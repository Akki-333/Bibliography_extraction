import streamlit as st
import pytesseract
from PIL import Image
import fitz  # PyMuPDF for PDF text extraction
import io
import docx
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from pptx import Presentation
from pptx.util import Inches
import pandas as pd
import re
import spacy

# Initialize spaCy model
nlp = spacy.load("en_core_web_sm")

# Define a simple summarization function (can be customized)
def simple_summarize(text):
    sentences = text.split('. ')
    summary = '. '.join(sentences[:2])  # Just take the first two sentences as a summary
    return summary

# Function to extract bibliography information using regex
def extract_bibliography_info(text):
    bibliography_info = {
        'title': '',
        'author': '',
        'year': ''
    }
    
    text = text.strip().replace('\n', ' ').replace('  ', ' ')
    
    # Improved regex patterns
    author_pattern = re.compile(r'(?:by\s+|written\s+by\s+|author\s+)\s*([A-Z][a-zA-Z\s,]+)', re.IGNORECASE)
    year_pattern = re.compile(r'\b(\d{4})\b')
    title_pattern = re.compile(r'\btitle\s*:\s*([A-Z][a-zA-Z\s]+(?:[a-zA-Z]))', re.IGNORECASE)
    
    # Attempt matching for author
    author_match = author_pattern.search(text)
    if author_match:
        bibliography_info['author'] = author_match.group(1).strip()
    
    # Attempt matching for year
    year_match = year_pattern.search(text)
    if year_match:
        bibliography_info['year'] = year_match.group(0).strip()
    
    # Attempt matching for title
    title_match = title_pattern.search(text)
    if title_match:
        bibliography_info['title'] = title_match.group(1).strip()

    # If no matches are found, attempt NLP processing for a better approach
    if not bibliography_info['author']:
        bibliography_info['author'] = extract_author_from_nlp(text)
    
    if not bibliography_info['year']:
        bibliography_info['year'] = extract_year_from_nlp(text)
    
    if not bibliography_info['title']:
        bibliography_info['title'] = extract_title_from_nlp(text)
    
    # Ensure title and author are not the same
    if bibliography_info['title'] == bibliography_info['author']:
        bibliography_info['title'] = extract_title_from_nlp(text)
    
    return bibliography_info

# Function to extract author using NLP
def extract_author_from_nlp(text):
    doc = nlp(text)
    for ent in doc.ents:
        if ent.label_ == "PERSON":
            return ent.text
    return ''

# Function to extract year using NLP
def extract_year_from_nlp(text):
    doc = nlp(text)
    for ent in doc.ents:
        if ent.label_ == "DATE":
            return ent.text[-4:]  # Extract year
    return ''

# Function to extract title using NLP
def extract_title_from_nlp(text):
    doc = nlp(text)
    sentences = text.split('. ')
    # Heuristic: First sentence is often the title
    if len(sentences) > 0:
        return sentences[0]
    return ''

# Function to extract text from an image
def extract_text_from_image(image_bytes):
    try:
        image = Image.open(io.BytesIO(image_bytes))
        text = pytesseract.image_to_string(image)
        return text
    except Exception as e:
        st.error(f"Error extracting text from image: {e}")
        return ""

# Function to extract text from a PDF using PyMuPDF
def extract_text_from_pdf(pdf_bytes):
    try:
        pdf_document = fitz.open(stream=pdf_bytes, filetype="pdf")
        text = ""
        for page_num in range(len(pdf_document)):
            page = pdf_document.load_page(page_num)
            text += page.get_text("text")
        return text
    except Exception as e:
        st.error(f"Error extracting text from PDF: {e}")
        return ""

# Function to format bibliography information for download
def format_bibliography_info(bibliography_info):
    return f"Title: {bibliography_info['title']}\n\nAuthor: {bibliography_info['author']}\n\nYear: {bibliography_info['year']}"

# Function to download as Word
def download_word(content):
    doc = docx.Document()
    doc.add_paragraph(content)
    output_file = io.BytesIO()
    doc.save(output_file)
    output_file.seek(0)
    return output_file

# Function to download as PDF
def download_pdf(content):
    output_file = io.BytesIO()
    c = canvas.Canvas(output_file, pagesize=letter)
    text_object = c.beginText(40, 750)
    text_object.setFont("Helvetica", 12)
    lines = content.splitlines()
    for line in lines:
        text_object.textLine(line)
    c.drawText(text_object)
    c.showPage()
    c.save()
    output_file.seek(0)
    return output_file

# Function to download as PowerPoint
def download_ppt(content):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    textbox = slide.shapes.add_textbox(left=Inches(1), top=Inches(1), width=Inches(8), height=Inches(5))
    text_frame = textbox.text_frame
    p = text_frame.add_paragraph()
    p.text = content
    output_file = io.BytesIO()
    prs.save(output_file)
    output_file.seek(0)
    return output_file

# Function to download as Excel
def download_excel(content):
    df = pd.DataFrame({"Content": [content]})
    output_file = io.BytesIO()
    df.to_excel(output_file, index=False)
    output_file.seek(0)
    return output_file

# Initialize session state variables
if 'page' not in st.session_state:
    st.session_state['page'] = 'main'

if 'file_name' not in st.session_state:
    st.session_state['file_name'] = 'bibliography_content'

if 'theme' not in st.session_state:
    st.session_state['theme'] = 'light'

# Apply custom theme based on session state
theme_styles = {
    'light': """
    <style>
    .main-container {
        background-color: #f5f5f5;
        padding: 20px;
        border-radius: 10px;
    }
    .btn {
        background-color: #4CAF50;
        color: white;
        padding: 10px 20px;
        border: none;
        border-radius: 5px;
        cursor: pointer;
        font-size: 16px;
    }
    .btn:hover {
        background-color: #45a049;
    }
    .header {
        font-size: 36px;
        font-weight: 600;
        color: #333;
    }
    .subheader {
        font-size: 24px;
        font-weight: 500;
        color: #555;
    }
    .text-box {
        border: 2px solid #ddd;
        padding: 10px;
        border-radius: 5px;
        background-color: #fff;
        color: black;
        max-height: 400px;
        overflow-y: auto;
        white-space: pre-wrap;
        font-family: monospace;
    }
    .file-name-container {
        display: flex;
        align-items: center;
        margin-bottom: 20px;
    }
    .file-name-container input {
        flex: 1;
        padding: 10px;
        border-radius: 5px;
        border: 1px solid #ccc;
    }
    .file-name-container button {
        margin-left: 10px;
    }
    </style>
    """,
    'dark': """
    <style>
    .main-container {
        background-color: #333;
        padding: 20px;
        border-radius: 10px;
        color: #f5f5f5;
    }
    .btn {
        background-color: #4CAF50;
        color: white;
        padding: 10px 20px;
        border: none;
        border-radius: 5px;
        cursor: pointer;
        font-size: 16px;
    }
    .btn:hover {
        background-color: #45a049;
    }
    .header {
        font-size: 36px;
        font-weight: 600;
        color: #f5f5f5;
    }
    .subheader {
        font-size: 24px;
        font-weight: 500;
        color: #ddd;
    }
    .text-box {
        border: 2px solid #555;
        padding: 10px;
        border-radius: 5px;
        background-color: #444;
        color: #f5f5f5;
        max-height: 400px;
        overflow-y: auto;
        white-space: pre-wrap;
        font-family: monospace;
    }
    .file-name-container {
        display: flex;
        align-items: center;
        margin-bottom: 20px;
    }
    .file-name-container input {
        flex: 1;
        padding: 10px;
        border-radius: 5px;
        border: 1px solid #777;
    }
    .file-name-container button {
        margin-left: 10px;
    }
    </style>
    """
}

# Apply the selected theme style
st.markdown(theme_styles[st.session_state['theme']], unsafe_allow_html=True)

# Main and About page content
def main_page():
    st.title("Bibliography Extraction Tool")
    
    menu = ["Extract Bibliography", "Summarize Text"]
    choice = st.sidebar.selectbox("Choose an action", menu)
    
    if choice == "Extract Bibliography":
        st.subheader("Upload a file to extract bibliography information")

        uploaded_file = st.file_uploader("Choose a file", type=["pdf", "png", "jpg", "jpeg"])

        if uploaded_file is not None:
            if uploaded_file.type == "application/pdf":
                text = extract_text_from_pdf(uploaded_file.read())
            elif uploaded_file.type in ["image/png", "image/jpeg", "image/jpg"]:
                text = extract_text_from_image(uploaded_file.read())
            else:
                st.error("Unsupported file type.")
                return

            if text:
                bibliography_info = extract_bibliography_info(text)
                formatted_bibliography = format_bibliography_info(bibliography_info)

                st.markdown("### Extracted Bibliography Information")
                st.text_area("Extracted Text", value=text, height=300)

                st.markdown("### Extracted Bibliography Details")
                st.write(f"*Title:* {bibliography_info['title']}") 
                st.write(f"*Author:* {bibliography_info['author']}") 
                st.write(f"*Year:* {bibliography_info['year']}")

                # Download options
                st.download_button("Download as Word", download_word(formatted_bibliography), file_name="bibliography.docx")
                st.download_button("Download as PDF", download_pdf(formatted_bibliography), file_name="bibliography.pdf")
                st.download_button("Download as PowerPoint", download_ppt(formatted_bibliography), file_name="bibliography.pptx")
                st.download_button("Download as Excel", download_excel(formatted_bibliography), file_name="bibliography.xlsx")

    elif choice == "Summarize Text":
        st.subheader("Upload a file to summarize the content")
        uploaded_file = st.file_uploader("Choose a file", type=["pdf", "png", "jpg", "jpeg"])

        if uploaded_file is not None:
            if uploaded_file.type == "application/pdf":
                text = extract_text_from_pdf(uploaded_file.read())
            elif uploaded_file.type in ["image/png", "image/jpeg", "image/jpg"]:
                text = extract_text_from_image(uploaded_file.read())
            else:
                st.error("Unsupported file type.")
                return

            if text:
                summary = simple_summarize(text)

                st.markdown("### Original Text")
                st.text_area("Original Text", value=text, height=300)

                st.markdown("### Summary")
                st.text_area("Summary", value=summary, height=150)

                # Download options
                st.download_button("Download Summary as Word", download_word(summary), file_name="summary.docx")
                st.download_button("Download Summary as PDF", download_pdf(summary), file_name="summary.pdf")
                st.download_button("Download Summary as PowerPoint", download_ppt(summary), file_name="summary.pptx")
                st.download_button("Download Summary as Excel", download_excel(summary), file_name="summary.xlsx")

def about_page():
    st.title("About")
    st.write("""  
        # Bibliography Extraction and Summarization Tool
        
        This tool allows users to upload various types of files (PDFs, images, etc.) to extract 
        key bibliographic details, such as title, author, and publication year, as well as generate summaries 
        of the content in these files. It utilizes advanced text extraction techniques including Optical 
        Character Recognition (OCR) for images and PyMuPDF for PDFs.
        
        ## Features:
        - *Extract Bibliography Information*: 
          Extract the title, author, and year from a file, whether it's an image, PDF, or other formats.
          The tool uses both traditional regex matching and Natural Language Processing (NLP) to improve accuracy.
          
        - *Summarize Text*: 
          Summarize the text content of any uploaded file, whether it's an image, PDF, or other types. 
          The tool uses simple summarization by extracting the first two sentences but can be customized 
          further for more detailed summaries.

        - *Download Options*: 
          After extracting bibliography information or summarizing content, the user can download the results 
          in different formats such as:
          - Word (DOCX)
          - PDF
          - PowerPoint (PPTX)
          - Excel (XLSX)

        ## Technologies Used:
        - *Streamlit* for creating the interactive web interface.
        - *spaCy* for Natural Language Processing to extract title, author, and year.
        - *PyMuPDF* for extracting text from PDF documents.
        - *Pytesseract* for OCR (Optical Character Recognition) to extract text from images.
        - *ReportLab* for generating PDFs.
        - *python-pptx* for creating PowerPoint presentations.
        - *Pandas* for handling and exporting data in Excel format.
    """)

# Sidebar navigation for pages
page = st.sidebar.radio("Select a Page", ("Main", "About"))

if page == "Main":
    main_page()
else:
    about_page()